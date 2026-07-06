"""
RAG 核心引擎 v4.2
==================
v3.0 能力 + v4.0/v4.1 增强：
  - 多轮对话上下文注入 (query/history)
  - 流式输出 (stream_query)
  - 标签/分类管理 (load_documents tags / query tag_filter)
  - 全文搜索+关键词高亮 (fulltext_search)
  - URL/网页摄取 (add_url)
  - 向量库版本快照/回滚 (VersionManager 集成)
  - RAG fallback 到纯 LLM (低分时)

v4.2 重构：
  - 提取 _resolve_sources() / _run_fallback() / _stream_fallback() / _split_scored()
  - query() 与 stream_query() 共享同一检索链路，行为一致
  - 减少约 60 行重复代码

兼容 LangChain 1.3.0 (LCEL 风格)
"""

import hashlib
import os
import re
import shutil
import warnings
from pathlib import Path
from typing import List, Optional, Tuple, Dict, Any, Iterable, Generator, Callable

from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_community.retrievers import BM25Retriever
from rank_bm25 import BM25Okapi
from langchain_ollama import ChatOllama
from langchain_core.prompts import PromptTemplate
from langchain_core.embeddings import Embeddings
from langchain_core.output_parsers import StrOutputParser
from langchain_core.language_models.chat_models import BaseChatModel

import config


# =============================================================================
# LLM 工厂:支持 Ollama / OpenAI 兼容(OpenAI / 通义 / 智谱 / DeepSeek / 月之暗面)
# =============================================================================
def _create_llm(
    provider: str,
    model: str,
    ollama_base_url: str,
    openai_api_key: str,
    openai_base_url: str,
    openai_model: str,
    temperature: float = 0.1,
) -> BaseChatModel:
    """
    根据 LLM_PROVIDER 构造 LangChain ChatModel 实例。

    :param provider: "ollama" | "openai"
    :param model: Ollama 模型名(provider=ollama 时使用)
    :param ollama_base_url: Ollama 服务地址
    :param openai_api_key: OpenAI 兼容 API Key(provider=openai 时使用)
    :param openai_base_url: OpenAI 兼容 base URL(可指向通义/智谱/DeepSeek 等)
    :param openai_model: OpenAI 兼容模型名
    :param temperature: 采样温度,默认 0.1(略低,保证稳定)
    :return: LangChain BaseChatModel 子类实例
    """
    provider = (provider or "ollama").lower()

    if provider == "ollama":
        return ChatOllama(
            model=model,
            base_url=ollama_base_url,
            temperature=temperature,
        )

    if provider == "openai":
        # 懒加载:仅在用户显式选择 openai 时才 import
        from langchain_openai import ChatOpenAI

        if not openai_api_key:
            raise ValueError(
                "LLM_PROVIDER=openai 时必须设置 OPENAI_API_KEY 环境变量。"
                " 详见 .env.example"
            )
        return ChatOpenAI(
            model=openai_model or model,
            api_key=openai_api_key,
            base_url=openai_base_url,
            temperature=temperature,
        )

    raise ValueError(
        f"不支持的 LLM_PROVIDER: {provider!r},可选值: ollama / openai"
    )

# v4.0 新增：URL 抓取与版本管理（包内相对导入）
from .web_loader import WebLoader, is_valid_url
from .version_manager import VersionManager

warnings.filterwarnings("ignore")


# =============================================================================
# 1. 自定义文档加载器
# =============================================================================

class PDFLoader:
    """PDF 加载器：优先 jin-doc-parser，回退 PyPDF"""

    def __init__(self, file_path: str):
        self.file_path = file_path

    def load(self) -> List[Document]:
        if config.PDF_PARSER == "jin-doc-parser":
            try:
                return self._load_with_jin()
            except Exception as e:
                print(f"[RAG] jin-doc-parser 失败 ({e})，回退到 PyPDF")

        try:
            from langchain_community.document_loaders import PyPDFLoader
            loader = PyPDFLoader(self.file_path)
            docs = loader.load()
            fname = Path(self.file_path).name
            for i, doc in enumerate(docs):
                doc.metadata.setdefault("source", fname)
                doc.metadata.setdefault("page", i)
                doc.metadata.setdefault("file_type", ".pdf")
            return docs
        except Exception as e:
            print(f"[RAG] ❌ PDF 加载失败 {self.file_path}: {e}")
            return []

    def _load_with_jin(self) -> List[Document]:
        try:
            from jin_document_parser import PDFParser
        except ImportError:
            raise ImportError("jin-doc-parser 未安装")
        parser = PDFParser()
        result = parser.parse(self.file_path)
        md_text = result.get("markdown", result.get("text", ""))
        if not md_text.strip():
            return []
        fname = Path(self.file_path).name
        return [Document(
            page_content=md_text,
            metadata={"source": fname, "page": 0, "file_type": ".pdf"}
        )]


class DOCXLoader:
    """自定义 Word 文档加载器（python-docx，仅支持 .docx OOXML 格式）"""

    def __init__(self, file_path: str):
        self.file_path = file_path

    def load(self) -> List[Document]:
        ext = Path(self.file_path).suffix.lower()
        # .doc 是旧版二进制格式，python-docx 只支持 .docx
        if ext == ".doc":
            print(f"[RAG] ⚠️ .doc 文件是旧版二进制格式，python-docx 无法读取，请用 Word 另存为 .docx 后再上传")
            return []
        try:
            from docx import Document as DocxDocument
            docx = DocxDocument(self.file_path)
            full_text = []
            for para in docx.paragraphs:
                if para.text.strip():
                    full_text.append(para.text.strip())
            for table in docx.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        full_text.append(" | ".join(row_text))
            if not full_text:
                return []
            fname = Path(self.file_path).name
            return [Document(
                page_content="\n".join(full_text),
                metadata={"source": fname, "page": 0, "file_type": ".docx"}
            )]
        except Exception as e:
            err_msg = str(e).lower()
            if "package" in err_msg or "not found" in err_msg or "olefile" in err_msg or "docx" in err_msg:
                print(f"[RAG] ⚠️ .doc 文件是旧版二进制格式，python-docx 无法读取，请用 Word 另存为 .docx 后再上传（{self.file_path}）")
            else:
                print(f"[RAG] ❌ DOCX 加载失败 {self.file_path}: {e}")
            return []


class ExcelLoader:
    """Excel 加载器（.xlsx 走 openpyxl，.xls 走 xlrd，无依赖时降级）"""

    def __init__(self, file_path: str):
        self.file_path = file_path

    def load(self) -> List[Document]:
        # 根据文件后缀选择加载路径，老格式 .xls 与新格式 .xlsx 走不同解析库
        ext = Path(self.file_path).suffix.lower()
        if ext == ".xls":
            return self._load_xls()
        return self._load_xlsx()

    def _load_xlsx(self) -> List[Document]:
        """使用 openpyxl 解析 .xlsx 文件，加载失败时返回空列表而非抛出异常。"""
        try:
            import openpyxl
            wb = openpyxl.load_workbook(self.file_path, read_only=True, data_only=True)
            fname = Path(self.file_path).name
            all_docs = []
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                rows_data = []
                for row in sheet.iter_rows(values_only=True):
                    if any(cell is not None for cell in row):
                        row_text = " | ".join(
                            str(cell) if cell is not None else "" for cell in row
                        )
                        rows_data.append(row_text)
                if rows_data:
                    content = f"[工作表: {sheet_name}]\n" + "\n".join(rows_data)
                    all_docs.append(Document(
                        page_content=content,
                        metadata={"source": fname, "page": 0, "file_type": ".xlsx", "sheet": sheet_name}
                    ))
            wb.close()
            return all_docs
        except Exception as e:
            print(f"[RAG] ❌ Excel (.xlsx) 加载失败 {self.file_path}: {e}")
            return []

    def _load_xls(self) -> List[Document]:
        """老格式 .xls 用 xlrd；不可用时降级为返回空列表并给出友好提示。"""
        try:
            import xlrd
        except ImportError:
            print("[RAG] ⚠️ .xls 格式需要 xlrd（pip install xlrd==2.0.1），已跳过")
            return []
        try:
            book = xlrd.open_workbook(self.file_path)
            fname = Path(self.file_path).name
            all_docs = []
            for sheet in book.sheets():
                rows_data = []
                for row_idx in range(sheet.nrows):
                    row = sheet.row_values(row_idx)
                    if any(c not in (None, "") for c in row):
                        rows_data.append(" | ".join(str(c) for c in row))
                if rows_data:
                    content = f"[工作表: {sheet.name}]\n" + "\n".join(rows_data)
                    all_docs.append(Document(
                        page_content=content,
                        metadata={"source": fname, "page": 0, "file_type": ".xls", "sheet": sheet.name}
                    ))
            return all_docs
        except Exception as e:
            print(f"[RAG] ❌ Excel (.xls) 加载失败 {self.file_path}: {e}")
            return []


class PPTXLoader:
    """自定义 PPT 加载器（python-pptx，仅支持 .pptx / .pptx OOXML 格式）"""

    def __init__(self, file_path: str):
        self.file_path = file_path

    def load(self) -> List[Document]:
        ext = Path(self.file_path).suffix.lower()
        # .ppt 是旧版二进制格式，python-pptx 完全不支持
        if ext == ".ppt":
            print(f"[RAG] ⚠️ 不支持 .ppt 旧格式（{self.file_path}），请另存为 .pptx 后重试")
            return []
        try:
            from pptx import Presentation
            prs = Presentation(self.file_path)
            fname = Path(self.file_path).name
            all_docs = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text and shape.text.strip():
                        slide_text.append(shape.text.strip())
                if slide_text:
                    content = "\n".join(slide_text)
                    all_docs.append(Document(
                        page_content=content,
                        metadata={"source": fname, "page": slide_num - 1, "file_type": ".pptx"}
                    ))
            return all_docs
        except Exception as e:
            err_msg = str(e).lower()
            if "package" in err_msg or "not found" in err_msg or "olefile" in err_msg:
                print(f"[RAG] ⚠️ .ppt 文件是旧版二进制格式，python-pptx 无法读取，请用 PowerPoint 另存为 .pptx 后再上传")
            else:
                print(f"[RAG] ❌ PPTX 加载失败 {self.file_path}: {e}")
            return []


class TextLoader:
    """通用文本加载器（TXT, MD, HTML, JSON, XML, CSV, RTF）"""

    def __init__(self, file_path: str):
        self.file_path = file_path

    def load(self) -> List[Document]:
        try:
            ext = Path(self.file_path).suffix.lower()
            fname = Path(self.file_path).name
            if ext == ".rtf":
                return self._load_rtf()
            for encoding in ["utf-8", "gbk", "gb2312", "latin-1"]:
                try:
                    with open(self.file_path, "r", encoding=encoding) as f:
                        text = f.read()
                    break
                except (UnicodeDecodeError, UnicodeError):
                    continue
            else:
                print(f"[RAG] ❌ 无法解码文件 {self.file_path}")
                return []
            if not text.strip():
                return []
            return [Document(
                page_content=text,
                metadata={"source": fname, "page": 0, "file_type": ext}
            )]
        except Exception as e:
            print(f"[RAG] ❌ 文本加载失败 {self.file_path}: {e}")
            return []

    def _load_rtf(self) -> List[Document]:
        try:
            from striprtf.striprtf import rtf_to_text
            with open(self.file_path, "r", encoding="utf-8", errors="ignore") as f:
                rtf_content = f.read()
            text = rtf_to_text(rtf_content)
        except ImportError:
            with open(self.file_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
            text = re.sub(r'\\[a-z]+\d*\s?', '', text)
            text = re.sub(r'[{}]', '', text)
        fname = Path(self.file_path).name
        if not text.strip():
            return []
        return [Document(
            page_content=text,
            metadata={"source": fname, "page": 0, "file_type": ".rtf"}
        )]


def get_loader(file_path: str):
    """
    根据文件扩展名返回对应的加载器实例。

    旧版二进制格式 .doc / .ppt 由专门类提示「需要另存为新格式」，
    避免上游日志出现「DOCX 加载失败」之类的误导信息。
    """
    ext = Path(file_path).suffix.lower()
    fname = Path(file_path).name

    # 旧版二进制格式 .doc / .ppt：python-docx / python-pptx 都不支持，
    # 直接返回 None 并在调用方给出统一提示，不再让 DOCXLoader 触发异常
    if ext == ".doc":
        print(f"[RAG] ⚠️ 「{fname}」是旧版 .doc 二进制格式，请用 Word 另存为 .docx 后再上传")
        return None
    if ext == ".ppt":
        print(f"[RAG] ⚠️ 「{fname}」是旧版 .ppt 二进制格式，请用 PowerPoint 另存为 .pptx 后再上传")
        return None

    loaders = {
        ".pdf": PDFLoader, ".docx": DOCXLoader,
        ".xlsx": ExcelLoader, ".xls": ExcelLoader,
        ".pptx": PPTXLoader,
    }
    text_exts = {".txt", ".md", ".html", ".htm", ".json", ".xml", ".csv", ".rtf"}
    if ext in loaders:
        return loaders[ext](file_path)
    elif ext in text_exts:
        return TextLoader(file_path)
    else:
        print(f"[RAG] ⚠️ 不支持的格式: {ext}")
        return None


# =============================================================================
# 2. 增强分片器
# =============================================================================

class EnhancedTextSplitter:
    """增强分片器：表格保护 + 标题绑定 + 去重 + 30% overlap"""

    def __init__(
        self,
        chunk_size: int = config.CHUNK_SIZE,
        chunk_overlap: int = config.CHUNK_OVERLAP,
        table_aware: bool = config.TABLE_AWARE,
        header_binding: bool = config.HEADER_BINDING,
    ):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.table_aware = table_aware
        self.header_binding = header_binding
        self._splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            separators=["\n\n", "\n", "。", "！", "？", "；", " ", ""],
            length_function=len,
        )

    def split_documents(self, docs: List[Document]) -> List[Document]:
        all_chunks = []
        for doc in docs:
            text = doc.page_content
            base_meta = doc.metadata.copy()
            tables = []
            if self.table_aware:
                text, tables = self._extract_tables(text)
            raw_chunks = self._splitter.create_documents([text], metadatas=[base_meta])
            if self.table_aware and tables:
                raw_chunks = self._restore_tables(raw_chunks, tables)
            if self.header_binding:
                raw_chunks = self._bind_headers(raw_chunks)
            all_chunks.extend(raw_chunks)
        all_chunks = self._deduplicate(all_chunks)
        print(f"[RAG] ✂️ 分片完成: {len(docs)} 个文档 → {len(all_chunks)} 个块")
        return all_chunks

    def _extract_tables(self, text: str) -> Tuple[str, List[str]]:
        tables = []
        table_pattern = r'(\|.+\|\n\|[-:\s|]+\|\n(?:\|.+\|\n?)*)'
        def replace_table(match):
            tables.append(match.group(1))
            return f"__TABLE_PLACEHOLDER_{len(tables) - 1}__"
        text = re.sub(table_pattern, replace_table, text)
        return text, tables

    def _restore_tables(self, chunks: List[Document], tables: List[str]) -> List[Document]:
        for chunk in chunks:
            for i, table in enumerate(tables):
                placeholder = f"__TABLE_PLACEHOLDER_{i}__"
                if placeholder in chunk.page_content:
                    chunk.page_content = chunk.page_content.replace(placeholder, table)
        return chunks

    def _bind_headers(self, chunks: List[Document]) -> List[Document]:
        if len(chunks) <= 1:
            return chunks
        merged = []
        i = 0
        while i < len(chunks):
            current = chunks[i]
            if (current.page_content.strip().startswith("#") and
                    len(current.page_content) < self.chunk_size // 2 and
                    i + 1 < len(chunks)):
                next_chunk = chunks[i + 1]
                merged_content = current.page_content + "\n\n" + next_chunk.page_content
                merged_meta = {**current.metadata, **next_chunk.metadata}
                merged_doc = Document(
                    page_content=merged_content[:self.chunk_size * 2],
                    metadata=merged_meta
                )
                merged.append(merged_doc)
                i += 2
            else:
                merged.append(current)
                i += 1
        return merged

    def _deduplicate(self, chunks: List[Document]) -> List[Document]:
        seen = set()
        unique = []
        for chunk in chunks:
            content_hash = hash(chunk.page_content.strip())
            if content_hash not in seen:
                seen.add(content_hash)
                unique.append(chunk)
        removed = len(chunks) - len(unique)
        if removed > 0:
            print(f"[RAG] ✂️ 去重: 移除 {removed} 个重复块")
        return unique


# =============================================================================
# 3. 双模式 Embedding
# =============================================================================

class HuggingFaceBGEEmbedding(Embeddings):
    """HuggingFace BGE 嵌入（离线加载）"""

    def __init__(self, model_path: str = config.BGE_EMBED_MODEL):
        self.model_path = model_path
        self._model = None
        self._load_model()

    def _load_model(self):
        try:
            from sentence_transformers import SentenceTransformer
            print(f"[RAG] 📐 加载 BGE 嵌入模型: {self.model_path}")
            self._model = SentenceTransformer(self.model_path)
            print(f"[RAG] 📐 BGE 嵌入模型加载完成")
        except Exception as e:
            print(f"[RAG] ❌ BGE 嵌入模型加载失败: {e}")
            raise

    def embed_documents(self, texts: List[str]) -> List[List[float]]:
        embeddings = self._model.encode(texts, normalize_embeddings=True)
        return embeddings.tolist()

    def embed_query(self, text: str) -> List[float]:
        embedding = self._model.encode([text], normalize_embeddings=True)
        return embedding[0].tolist()


class OllamaEmbeddingWrapper(Embeddings):
    """Ollama 嵌入包装器"""

    def __init__(self, model: str = config.OLLAMA_EMBED_MODEL,
                 base_url: str = config.OLLAMA_BASE_URL):
        self.model = model
        self.base_url = base_url

    def embed_documents(self, texts: List[str]) -> List[List[float]]:
        import ollama
        client = ollama.Client(host=self.base_url)
        embeddings = []
        for text in texts:
            resp = client.embeddings(model=self.model, prompt=text)
            embeddings.append(resp["embedding"])
        return embeddings

    def embed_query(self, text: str) -> List[float]:
        import ollama
        client = ollama.Client(host=self.base_url)
        resp = client.embeddings(model=self.model, prompt=text)
        return resp["embedding"]


def create_embedding() -> Embeddings:
    if config.EMBEDDING_PROVIDER == "huggingface":
        return HuggingFaceBGEEmbedding()
    elif config.EMBEDDING_PROVIDER == "ollama":
        return OllamaEmbeddingWrapper()
    else:
        raise ValueError(f"未知的 Embedding Provider: {config.EMBEDDING_PROVIDER}")


# =============================================================================
# 4. CrossEncoder Reranker
# =============================================================================

class CrossEncoderReranker:
    """CrossEncoder 精排器"""

    def __init__(
        self,
        model_path: str = config.RERANK_MODEL,
        initial_k: int = config.RERANK_INITIAL_K,
        final_k: int = config.RERANK_FINAL_K,
        score_threshold: float = config.RERANK_SCORE_THRESHOLD,
    ):
        self.model_path = model_path
        self.initial_k = initial_k
        self.final_k = final_k
        self.score_threshold = score_threshold
        self._model = None
        if config.RERANK_ENABLED:
            self._load_model()

    def _load_model(self):
        try:
            from sentence_transformers import CrossEncoder
            print(f"[RAG] 🔄 加载 Rerank 模型: {self.model_path}")
            self._model = CrossEncoder(self.model_path)
            print(f"[RAG] 🔄 Rerank 模型加载完成")
        except Exception as e:
            print(f"[RAG] ❌ Rerank 模型加载失败: {e}，将跳过 Rerank")
            config.RERANK_ENABLED = False

    def rerank(self, query: str, docs: List[Document]) -> List[Tuple[Document, float]]:
        if not config.RERANK_ENABLED or self._model is None:
            return [(doc, 0.0) for doc in docs]
        if not docs:
            return []
        pairs = [(query, doc.page_content) for doc in docs]
        scores = self._model.predict(pairs)
        scored_docs = list(zip(docs, scores.tolist()))
        scored_docs.sort(key=lambda x: x[1], reverse=True)
        if self.score_threshold > 0:
            scored_docs = [(doc, score) for doc, score in scored_docs
                           if score >= self.score_threshold]
        scored_docs = scored_docs[:self.final_k]
        return scored_docs


# =============================================================================
# 5. RRF 混合检索辅助函数
# =============================================================================

def _doc_id(doc: Document) -> str:
    """生成文档唯一 ID：来源 + 内容 hash。"""
    source = doc.metadata.get("source", "")
    content_hash = hash(doc.page_content.strip())
    return f"{source}_{content_hash}"


# =============================================================================
# 6. RAG 引擎主类
# =============================================================================

class RAGEngine:
    """RAG 核心引擎（LCEL 风格，兼容 LangChain 1.3.0）"""

    def __init__(self):
        # 向量库与文档状态
        self.vectorstore: Optional[FAISS] = None
        self.doc_names: List[str] = []
        self.chunks: List[Document] = []
        self.bm25_retriever: Optional[BM25Retriever] = None
        # 模型与处理组件
        self.reranker: Optional[CrossEncoderReranker] = None
        self.embedding: Optional[Embeddings] = None
        self.llm: Optional[ChatOllama] = None
        self._splitter: Optional[EnhancedTextSplitter] = None
        # v4.0 新增：版本管理器（延迟初始化，配置关闭时为 None）
        self.version_manager: Optional[VersionManager] = None
        self._init_components()

    def _init_components(self):
        """初始化 Embedding / LLM / Reranker / 分片器 / 版本管理"""
        print("[RAG] 🚀 初始化 RAG 引擎...")
        try:
            self.embedding = create_embedding()
        except Exception as e:
            print(f"[RAG] ❌ 嵌入模型初始化失败: {e}")
            self.embedding = None

        try:
            self.llm = _create_llm(
                provider=config.LLM_PROVIDER,
                model=config.LLM_MODEL,
                ollama_base_url=config.OLLAMA_BASE_URL,
                openai_api_key=config.OPENAI_API_KEY,
                openai_base_url=config.OPENAI_BASE_URL,
                openai_model=config.OPENAI_MODEL,
            )
            print(f"[RAG] 🤖 LLM 已连接: provider={config.LLM_PROVIDER} model={self.llm.model if hasattr(self.llm, 'model') else config.LLM_MODEL}")
        except Exception as e:
            print(f"[RAG] ❌ LLM 连接失败: {e}")
            self.llm = None

        self.reranker = CrossEncoderReranker()
        self._splitter = EnhancedTextSplitter()
        # 版本管理器：配置关闭时不创建
        if getattr(config, "VERSION_ENABLED", True):
            try:
                self.version_manager = VersionManager(
                    vector_store_path=config.VECTOR_STORE_PATH,
                )
                print(f"[RAG] 📦 版本管理已启用: {self.version_manager.snapshot_dir}")
            except Exception as e:
                print(f"[RAG] ⚠️ 版本管理器初始化失败: {e}")
                self.version_manager = None
        self.restore_vectorstore()
        print("[RAG] 🚀 RAG 引擎初始化完成")

    def load_documents(
        self,
        file_paths: Optional[List[str]] = None,
        urls: Optional[List[str]] = None,
        tags: Optional[List[str]] = None,
        auto_snapshot: bool = True,
    ) -> Tuple[bool, str]:
        """
        加载文档到向量库（支持本地文件 + URL + 标签）

        Args:
            file_paths: 本地文件路径列表
            urls: 网页 URL 列表（用 WebLoader 抓取）
            tags: 标签列表，写入每个 chunk 的 metadata['tags']
            auto_snapshot: 加载前是否自动创建版本快照

        Returns:
            (是否成功, 描述信息)
        """
        # 合并 file_paths 和 urls，统一处理
        all_files: List[str] = list(file_paths or [])
        if urls:
            for url in urls:
                if not is_valid_url(url):
                    print(f"[RAG] ⚠️ URL 格式不合法: {url}")
                    continue
                loader = WebLoader(url)
                docs = loader.load()
                if not docs:
                    continue
                # 用 URL 的 md5 摘要作为文件名，避免特殊字符问题
                url_hash = hashlib.md5(url.encode("utf-8")).hexdigest()[:16]
                url_name = f"url_{url_hash}.html"
                tmp_path = os.path.join(config.UPLOADED_DOCS_PATH, url_name)
                os.makedirs(config.UPLOADED_DOCS_PATH, exist_ok=True)
                with open(tmp_path, "w", encoding="utf-8") as f:
                    f.write(docs[0].page_content)
                all_files.append(tmp_path)

        if not all_files:
            return False, "未提供文件或 URL"
        if self.embedding is None:
            return False, "嵌入模型未初始化"

        # 自动快照：变更前保护现场
        if auto_snapshot and self.vectorstore is not None and self.version_manager:
            try:
                self.version_manager.snapshot(
                    label="auto_before_load",
                    note=f"加载 {len(all_files)} 个文件前自动快照",
                    doc_count=len(self.doc_names),
                )
            except Exception as e:
                print(f"[RAG] ⚠️ 自动快照失败（不影响加载）: {e}")

        all_new_chunks: List[Document] = []
        loaded_names: List[str] = []
        # 标签统一成 list 写进 metadata
        tag_list = list(tags) if tags else []

        for fp in all_files:
            if not os.path.exists(fp):
                print(f"[RAG] ⚠️ 文件不存在: {fp}")
                continue
            size_mb = os.path.getsize(fp) / (1024 * 1024)
            if size_mb > config.MAX_FILE_SIZE_MB:
                print(f"[RAG] ⚠️ 文件过大 ({size_mb:.1f}MB): {fp}")
                continue
            loader = get_loader(fp)
            if loader is None:
                continue
            try:
                docs = loader.load()
                if not docs:
                    print(f"[RAG] ⚠️ 未提取到内容: {fp}")
                    continue
                if config.CHUNK_STRATEGY == "enhanced":
                    chunks = self._splitter.split_documents(docs)
                else:
                    base_splitter = RecursiveCharacterTextSplitter(
                        chunk_size=config.CHUNK_SIZE,
                        chunk_overlap=config.CHUNK_OVERLAP,
                    )
                    chunks = base_splitter.split_documents(docs)
                # 写入标签到每个 chunk 的 metadata
                if tag_list:
                    for c in chunks:
                        c.metadata["tags"] = tag_list
                all_new_chunks.extend(chunks)
                loaded_names.append(Path(fp).name)
            except Exception as e:
                print(f"[RAG] ❌ 加载失败 {fp}: {e}")
                continue

        if not all_new_chunks:
            return False, "没有成功加载任何文档内容"

        try:
            if self.vectorstore is None:
                self.vectorstore = FAISS.from_documents(all_new_chunks, self.embedding)
            else:
                self.vectorstore.add_documents(all_new_chunks)
            os.makedirs(config.VECTOR_STORE_PATH, exist_ok=True)
            self.vectorstore.save_local(config.VECTOR_STORE_PATH)
            for name in loaded_names:
                if name not in self.doc_names:
                    self.doc_names.append(name)
            self._rebuild_bm25()
            self.chunks.extend(all_new_chunks)
            msg = f"成功加载 {len(loaded_names)} 个文档，生成 {len(all_new_chunks)} 个文本块"
            print(f"[RAG] ✅ {msg}")
            return True, msg
        except Exception as e:
            return False, f"向量库构建失败: {e}"

    def add_url(self, url: str, tags: Optional[List[str]] = None) -> Tuple[bool, str]:
        """
        从 URL 摄取网页内容到知识库

        Args:
            url: 网页 URL
            tags: 关联标签

        Returns:
            (是否成功, 描述信息)
        """
        return self.load_documents(urls=[url], tags=tags)

    def query(
        self,
        question: str,
        history: Optional[List[Tuple[str, str]]] = None,
        tag_filter: Optional[List[str]] = None,
    ) -> Dict[str, Any]:
        """
        执行 RAG 查询（非流式）。

        Args:
            question: 用户问题
            history: 多轮对话历史，格式 [("user_q", "assistant_a"), ...]
            tag_filter: 仅在含这些标签的 chunk 中检索（AND 语义，None 不过滤）

        Returns:
            {"answer": str, "sources": List[(Document, float)], "fallback": bool (可选)}
        """
        # Step 1: 公共检索链路（KB 校验 + 混合检索 + 精排 + 阈值）
        status, scored_docs = self._resolve_sources(question, tag_filter)
        if status == "empty":
            return {"answer": "未找到相关内容", "sources": []}
        if status == "fallback":
            return self._run_fallback(question)

        # Step 2: 组装 LLM 输入
        top_docs, top_scores = self._split_scored(scored_docs)
        context = self._format_context(top_docs)
        history_text = self._format_history(history)
        full_question = (history_text + "\n" + question).strip() if history_text else question

        # Step 3: 调用 LLM 生成答案
        try:
            prompt = PromptTemplate(
                template=config.RAG_SYSTEM_PROMPT,
                input_variables=["context", "question"],
            )
            chain = prompt | self.llm | StrOutputParser()
            answer = chain.invoke({"context": context, "question": full_question})
        except Exception as e:
            return {"answer": f"⚠️ LLM 调用失败: {e}", "sources": []}

        return {"answer": answer, "sources": list(zip(top_docs, top_scores))}

    def stream_query(
        self,
        question: str,
        history: Optional[List[Tuple[str, str]]] = None,
        tag_filter: Optional[List[str]] = None,
    ) -> Generator[Dict[str, Any], None, None]:
        """
        流式 RAG 查询：逐 token yield 答案片段。

        与 query() 区别：检索/精排/prompt 组装在 yield 之间穿插进行，
        调用方拿到的是字典（含 done/sources 标志），便于前端逐步渲染。

        Yields:
            字典，type 字段取值:
              - "sources": 携带 sources 列表（先 yield 一次让前端显示来源）
              - "token":  单个答案片段
              - "done":   流结束信号
        """
        # 公共前处理：KB 校验
        if self.vectorstore is None:
            yield {"type": "token", "content": "⚠️ 知识库为空，请先上传文档"}
            yield {"type": "done", "content": ""}
            return
        if self.llm is None:
            yield {"type": "token", "content": "⚠️ LLM 未连接，请检查 Ollama 服务"}
            yield {"type": "done", "content": ""}
            return

        # Step 1: 公共检索链路
        status, scored_docs = self._resolve_sources(question, tag_filter)
        if status == "empty":
            yield {"type": "token", "content": "未找到相关内容"}
            yield {"type": "done", "content": ""}
            return
        if status == "fallback":
            # 先通知前端是 fallback 模式 + 留出 sources 面板位置
            yield {"type": "sources", "sources": [], "fallback": True}
            yield {"type": "token", "content": "🤖 "}
            for chunk in self._stream_fallback(question):
                yield chunk
            yield {"type": "done", "content": ""}
            return

        # Step 2: 先把 sources 推给前端，再启动 LLM 流
        top_docs, top_scores = self._split_scored(scored_docs)
        yield {"type": "sources", "sources": list(zip(top_docs, top_scores))}

        # Step 3: 组装 prompt 并启动流
        context = self._format_context(top_docs)
        history_text = self._format_history(history)
        full_question = (history_text + "\n" + question).strip() if history_text else question

        prompt = PromptTemplate(
            template=config.RAG_SYSTEM_PROMPT,
            input_variables=["context", "question"],
        )
        chain = prompt | self.llm | StrOutputParser()
        try:
            for chunk in chain.stream({"context": context, "question": full_question}):
                if chunk:
                    yield {"type": "token", "content": chunk}
        except Exception as e:
            yield {"type": "token", "content": f"\n⚠️ LLM 流式调用失败: {e}"}
        yield {"type": "done", "content": ""}

    # ------------------------------------------------------------------ #
    # 内部辅助方法（v4.2 提取：query/stream_query 共享检索链路）
    # ------------------------------------------------------------------ #
    def _resolve_sources(
        self,
        question: str,
        tag_filter: Optional[List[str]] = None,
    ) -> Tuple[str, Optional[List[Tuple[Document, float]]]]:
        """
        公共检索链路：KB 校验 → 混合检索 → Rerank → 阈值过滤 → fallback 决策。

        Returns:
            ("ok", scored_docs)  - 正常返回候选文档
            ("empty", None)      - KB 为空 / 检索为空 / 精排分数过低且未启用 fallback
            ("fallback", None)   - 精排后无结果，但 fallback 已启用（调用方走 LLM 模式）
        """
        if self.vectorstore is None or self.llm is None:
            return ("empty", None)

        retrieved_docs = self._hybrid_search(question, tag_filter)
        if not retrieved_docs:
            return ("empty", None)

        scored_docs = self.reranker.rerank(question, retrieved_docs)

        # 阈值过滤：使用 RAG_FALLBACK_THRESHOLD 作为「最低通过分」，
        # 避免"勉强相关"的低分文档（如 0.05~0.5）污染 LLM，导致答非所问
        # （v4.3 修复: 之前 RERANK_SCORE_THRESHOLD=0.0 时完全不过滤，
        #  RAG_FALLBACK_THRESHOLD 也只是装饰，没有真正生效）
        if scored_docs:
            max_score = max(s for _, s in scored_docs)
            if max_score < config.RAG_FALLBACK_THRESHOLD:
                if config.RAG_FALLBACK_ENABLED:
                    print(
                        f"[RAG Debug] 最高分 {max_score:.3f} < 阈值 "
                        f"{config.RAG_FALLBACK_THRESHOLD}，启用 LLM fallback"
                    )
                    return ("fallback", None)
                return ("empty", None)

        # 双保险：RERANK_SCORE_THRESHOLD 作为单文档最低分过滤线
        if config.RERANK_SCORE_THRESHOLD > 0:
            scored_docs = [
                (d, s) for d, s in scored_docs
                if s >= config.RERANK_SCORE_THRESHOLD
            ]

        if not scored_docs:
            if config.RAG_FALLBACK_ENABLED:
                print("[RAG Debug] 精排分数过低，启用 LLM fallback")
                return ("fallback", None)
            return ("empty", None)

        # 调试日志：检索 → 精排 链路
        print(f"[RAG Debug] query='{question[:50]}'")
        print(f"[RAG Debug] 混合检索 → {len(retrieved_docs)} | Rerank 后 → {len(scored_docs)}")
        for i, (doc, score) in enumerate(scored_docs[:3], 1):
            print(f"[RAG Debug]  Top{i}: {doc.metadata.get('source', '?')} score={score:.3f}")
        return ("ok", scored_docs)

    @staticmethod
    def _split_scored(
        scored_docs: List[Tuple[Document, float]],
    ) -> Tuple[List[Document], List[float]]:
        """把 (doc, score) 列表拆成两份平行列表，方便后续处理。"""
        top_docs = [d for d, _ in scored_docs]
        top_scores = [s for _, s in scored_docs]
        return top_docs, top_scores

    def _run_fallback(self, question: str) -> Dict[str, Any]:
        """执行非流式 fallback：直接用 LLM 回答（无上下文）。"""
        try:
            fallback_prompt = PromptTemplate(
                template=config.RAG_FALLBACK_PROMPT,
                input_variables=["question"],
            )
            chain = fallback_prompt | self.llm | StrOutputParser()
            answer = chain.invoke({"question": question})
            return {"answer": f"🤖 {answer}", "sources": [], "fallback": True}
        except Exception as e:
            return {"answer": f"⚠️ LLM fallback 失败: {e}", "sources": [], "fallback": True}

    def _stream_fallback(self, question: str) -> Generator[Dict[str, Any], None, None]:
        """执行流式 fallback：逐 token yield LLM 输出（无上下文）。"""
        try:
            fallback_prompt = PromptTemplate(
                template=config.RAG_FALLBACK_PROMPT,
                input_variables=["question"],
            )
            chain = fallback_prompt | self.llm | StrOutputParser()
            for chunk in chain.stream({"question": question}):
                if chunk:
                    yield {"type": "token", "content": chunk}
        except Exception as e:
            yield {"type": "token", "content": f"⚠️ LLM fallback 失败: {e}"}

    def fulltext_search(
        self,
        query: str,
        top_k: int = 20,
        highlight: bool = True,
        tag_filter: Optional[List[str]] = None,
    ) -> List[Dict[str, Any]]:
        """
        全文搜索（基于 BM25 关键词匹配，返回带高亮的 chunk 列表）

        Args:
            query: 搜索关键词
            top_k: 返回数量
            highlight: 是否返回带 <mark> 高亮标签的 preview
            tag_filter: 仅在含这些标签的 chunk 中检索

        Returns:
            列表，每项 {source, page, score, preview, content}
        """
        if not self.bm25_retriever or not self.chunks:
            return []
        # 标签过滤：先在 self.chunks 上筛，再传给 BM25
        chunks_pool = self.chunks
        if tag_filter:
            chunks_pool = [
                c for c in self.chunks
                if all(t in c.metadata.get("tags", []) for t in tag_filter)
            ]
            if not chunks_pool:
                return []

        # 重建一个针对 chunks_pool 的临时 BM25
        try:
            import jieba
            tokenized = [list(jieba.cut(c.page_content)) for c in chunks_pool]
            bm25 = BM25Okapi(tokenized)
            query_tokens = list(jieba.cut(query))
            scores = bm25.get_scores(query_tokens)
        except Exception as e:
            print(f"[RAG] ⚠️ 全文搜索失败: {e}")
            return []

        # 取 top_k 并组装结果
        ranked = sorted(
            zip(chunks_pool, scores), key=lambda x: x[1], reverse=True
        )[:top_k]
        results: List[Dict[str, Any]] = []
        for doc, score in ranked:
            if score <= 0:
                continue
            content = doc.page_content
            preview = content[:200].replace("\n", " ")
            if highlight and query:
                # 简单高亮：仅命中 query 第一个非空 token
                first_tok = next((t for t in query.split() if t.strip()), "")
                if first_tok and first_tok in preview:
                    preview = preview.replace(
                        first_tok, f"<mark>{first_tok}</mark>"
                    )
            results.append({
                "source": doc.metadata.get("source", "未知"),
                "page": doc.metadata.get("page", 0),
                "score": round(float(score), 3),
                "preview": preview,
                "content": content,
            })
        return results

    # ------------------------------------------------------------------ #
    # 内部辅助方法（v4.0 提取，避免 query/stream_query 重复）
    # ------------------------------------------------------------------ #
    def _hybrid_search(
        self,
        question: str,
        tag_filter: Optional[List[str]] = None,
    ) -> List[Document]:
        """执行 BM25 + 向量 RRF 混合检索，可选标签过滤。"""
        # 若有标签过滤，先把 chunks 池缩到符合标签的子集
        if tag_filter:
            filtered_chunks = [
                c for c in self.chunks
                if all(t in c.metadata.get("tags", []) for t in tag_filter)
            ]
            if not filtered_chunks:
                return []
            # 临时构造一个仅含 filtered_chunks 的 FAISS / BM25 子检索
            vector_docs = self._vector_search_filtered(question, filtered_chunks)
            bm25_docs = self._bm25_search_filtered(question, filtered_chunks)
        else:
            vector_docs = self._vector_search(question)
            bm25_docs = self._bm25_search(question)

        # RRF 融合
        rrf_k = 60
        doc_scores: Dict[str, float] = {}
        doc_map: Dict[str, Document] = {}
        for rank, (doc, score) in enumerate(vector_docs):
            doc_id = _doc_id(doc)
            doc_map[doc_id] = doc
            doc_scores[doc_id] = doc_scores.get(doc_id, 0.0) + \
                config.VECTOR_WEIGHT * (1.0 / (rrf_k + rank + 1))
        for rank, doc in enumerate(bm25_docs):
            doc_id = _doc_id(doc)
            doc_map[doc_id] = doc
            doc_scores[doc_id] = doc_scores.get(doc_id, 0.0) + \
                config.BM25_WEIGHT * (1.0 / (rrf_k + rank + 1))
        sorted_ids = sorted(doc_scores.keys(), key=lambda x: doc_scores[x], reverse=True)
        return [doc_map[did] for did in sorted_ids[:config.RETRIEVAL_INITIAL_K]]

    def _vector_search(self, question: str) -> List[Tuple[Document, float]]:
        """向量检索（FAISS similarity_search_with_score + 归一化）"""
        if self.vectorstore is None:
            return []
        try:
            k = config.RETRIEVAL_INITIAL_K
            results = self.vectorstore.similarity_search_with_score(question, k=k)
            if not results:
                return []
            raw = [s for _, s in results]
            mx, mn = max(raw), min(raw)
            rng = mx - mn if mx > mn else 1.0
            return [
                (doc, 1.0 - (s - mn) / rng) for doc, s in results
            ]
        except Exception as e:
            print(f"[RAG] ⚠️ 向量检索失败: {e}")
            return []

    def _bm25_search(self, question: str) -> List[Document]:
        """BM25 关键词检索（仅前 k 个）"""
        if not self.bm25_retriever:
            return []
        try:
            # LangChain 1.x：使用 invoke() 替代已删除的 get_relevant_documents()
            return self.bm25_retriever.invoke(question)[:config.RETRIEVAL_INITIAL_K]
        except Exception as e:
            print(f"[RAG] ⚠️ BM25 检索失败: {e}")
            return []

    def _vector_search_filtered(
        self, question: str, chunks: List[Document]
    ) -> List[Tuple[Document, float]]:
        """向量检索（限定 chunks 池；不做 FAISS 子索引重建，直接对 chunks 重嵌入会有性能损耗）。
        为简化实现，此处退化为：在主向量库中检索后按 chunk 身份映射回 filtered 集合。"""
        if self.vectorstore is None:
            return []
        try:
            k = config.RETRIEVAL_INITIAL_K
            results = self.vectorstore.similarity_search_with_score(question, k=k * 2)
            allowed_ids = {
                _doc_id(c) for c in chunks
            }
            raw_scores = [s for d, s in results if _doc_id(d) in allowed_ids]
            if not raw_scores:
                return []
            mx, mn = max(raw_scores), min(raw_scores)
            rng = mx - mn if mx > mn else 1.0
            return [
                (d, 1.0 - (s - mn) / rng)
                for d, s in results
                if _doc_id(d) in allowed_ids
            ]
        except Exception as e:
            print(f"[RAG] ⚠️ 向量检索（标签过滤）失败: {e}")
            return []

    def _bm25_search_filtered(
        self, question: str, chunks: List[Document]
    ) -> List[Document]:
        """BM25 在 chunks 池上临时构建索引。"""
        try:
            import jieba
            tokenized = [list(jieba.cut(c.page_content)) for c in chunks]
            bm25 = BM25Okapi(tokenized)
            scores = bm25.get_scores(list(jieba.cut(question)))
            ranked = sorted(zip(chunks, scores), key=lambda x: x[1], reverse=True)
            return [d for d, s in ranked[:config.RETRIEVAL_INITIAL_K] if s > 0]
        except Exception as e:
            print(f"[RAG] ⚠️ BM25（标签过滤）失败: {e}")
            return []

    def _format_context(self, top_docs: List[Document]) -> str:
        """把 top_docs 拼装成 LLM 上下文（来源标注）"""
        parts: List[str] = []
        for i, doc in enumerate(top_docs):
            src = doc.metadata.get("source", "未知")
            page = doc.metadata.get("page", "")
            tag_info = doc.metadata.get("tags", [])
            tag_str = f" [tags: {','.join(tag_info)}]" if tag_info else ""
            src_info = f"[来源{i+1}: {src}"
            if page != "" and isinstance(page, int):
                src_info += f" 第{page+1}页"
            src_info += f"{tag_str}]"
            parts.append(f"{src_info}\n{doc.page_content}")
        return "\n\n---\n\n".join(parts)

    def _format_history(self, history: Optional[List[Tuple[str, str]]]) -> str:
        """把多轮历史拼成对话串（截断过长内容以避免 token 爆炸）"""
        if not history:
            return ""
        max_turns = getattr(config, "MAX_HISTORY_TURNS", 5)
        max_chars = getattr(config, "MAX_HISTORY_CHARS", 1500)
        recent = history[-max_turns:]
        lines: List[str] = []
        total = 0
        for q, a in recent:
            line = f"用户: {q}\n助手: {a[:300]}"
            total += len(line)
            if total > max_chars:
                break
            lines.append(line)
        if not lines:
            return ""
        return "以下为之前的对话历史（可能相关）：\n" + "\n\n".join(lines)

    def _rebuild_bm25(self):
        """重建 BM25 索引（支持中文分词）"""
        if not config.BM25_ENABLED:
            self.bm25_retriever = None
            return

        if not self.chunks:
            self.bm25_retriever = None
            return

        try:
            import jieba

            # 中文分词后构建 BM25
            tokenized = []
            for chunk in self.chunks:
                words = list(jieba.cut(chunk.page_content))
                tokenized.append(words)

            self.bm25_retriever = BM25Retriever.from_documents(
                self.chunks,
                bm25_klass=BM25Okapi,
                preprocess_func=lambda text: list(jieba.cut(text)),
            )
            print(f"[RAG] 📊 BM25 索引已重建 ({len(self.chunks)} 个块)")
        except Exception as e:
            print(f"[RAG] ⚠️ BM25 索引重建失败: {e}")
            self.bm25_retriever = None

    def delete_document(self, doc_name: str, auto_snapshot: bool = True) -> bool:
        """
        删除单个文档的向量数据（精确删除，不影响其他文档）

        Args:
            doc_name: 文件名
            auto_snapshot: 删除前是否自动快照（防止误删可回滚）

        Returns:
            是否成功删除
        """
        if doc_name not in self.doc_names:
            print(f"[RAG] ⚠️ 文档不在知识库中: {doc_name}")
            return False

        if self.vectorstore is None:
            print("[RAG] ⚠️ 向量库为空，无需删除")
            return False

        # 删除前自动快照，便于误删回滚
        if auto_snapshot and self.version_manager:
            try:
                self.version_manager.snapshot(
                    label="auto_before_delete",
                    note=f"删除 {doc_name} 前自动快照",
                    doc_count=len(self.doc_names),
                )
            except Exception as e:
                print(f"[RAG] ⚠️ 删除前自动快照失败: {e}")

        try:
            # 找到该文档的所有向量块 ID
            docstore = self.vectorstore.docstore._dict
            ids_to_delete: List[str] = []
            for doc_id, doc in docstore.items():
                src = doc.metadata.get("source", "")
                if src == doc_name:
                    ids_to_delete.append(doc_id)

            if not ids_to_delete:
                print(f"[RAG] ⚠️ 向量库中未找到文档 {doc_name} 的数据")
                if doc_name in self.doc_names:
                    self.doc_names.remove(doc_name)
                return False

            # 从向量库中删除这些块
            self.vectorstore.delete(ids_to_delete)
            print(f"[RAG] ✓ 已从向量库删除 {len(ids_to_delete)} 个块 ({doc_name})")

            if doc_name in self.doc_names:
                self.doc_names.remove(doc_name)
            self.chunks = [c for c in self.chunks
                          if c.metadata.get("source", "") != doc_name]

            os.makedirs(config.VECTOR_STORE_PATH, exist_ok=True)
            self.vectorstore.save_local(config.VECTOR_STORE_PATH)
            self._rebuild_bm25()
            print(f"[RAG] ✓ BM25 索引已重建 (剩余 {len(self.chunks)} 个块)")
            print(f"[RAG] ✓ 向量库已更新保存")
            return True
        except Exception as e:
            print(f"[RAG] ❌ 删除文档失败: {e}")
            return False

    def restore_vectorstore(self, version_id: Optional[str] = None) -> None:
        """
        从磁盘恢复向量库

        Args:
            version_id: 指定要回滚到的快照 ID；为 None 时恢复主向量库
        """
        # 如果指定了 version_id，先回滚文件再正常加载
        if version_id and self.version_manager:
            if not self.version_manager.rollback(version_id):
                print(f"[RAG] ⚠️ 回滚到版本 {version_id} 失败，将尝试恢复主向量库")

        if not os.path.exists(config.VECTOR_STORE_PATH):
            print("[RAG] 📭 未找到已有向量库，将创建新的")
            return

        index_file = os.path.join(config.VECTOR_STORE_PATH, "index.faiss")
        if not os.path.exists(index_file):
            print("[RAG] 📭 向量库索引文件不存在，将创建新的")
            return

        if self.embedding is None:
            print("[RAG] ⚠️ 嵌入模型未初始化，跳过向量库恢复")
            return

        try:
            self.vectorstore = FAISS.load_local(
                config.VECTOR_STORE_PATH,
                self.embedding,
                allow_dangerous_deserialization=True,
            )
            print(f"[RAG] 📂 向量库已恢复: {self.vectorstore.index.ntotal} 个向量")

            seen = set()
            self.doc_names = []
            all_docs = self.vectorstore.docstore._dict
            for doc in all_docs.values():
                src = doc.metadata.get("source", "")
                if src and src not in seen:
                    seen.add(src)
                    self.doc_names.append(src)
            self.chunks = list(all_docs.values())
            self._rebuild_bm25()
            print(f"[RAG] 📂 恢复完成: {len(self.doc_names)} 个文档, {len(self.chunks)} 个块")
        except Exception as e:
            print(f"[RAG] ❌ 向量库恢复失败: {e}")
            self.vectorstore = None

    def clear_vectorstore(self):
        """清空整个向量库（仅清除向量数据，不删除原始文件）"""
        self.vectorstore = None
        self.doc_names = []
        self.chunks = []
        self.bm25_retriever = None

        # 删除向量库目录（仅向量数据，不碰 uploaded_docs）
        if os.path.exists(config.VECTOR_STORE_PATH):
            shutil.rmtree(config.VECTOR_STORE_PATH)
            os.makedirs(config.VECTOR_STORE_PATH, exist_ok=True)
            print("[RAG] ✓ 向量库已清空")
        else:
            print("[RAG] 向量库目录不存在，无需清空")

    def get_doc_names(self) -> List[str]:
        """返回知识库中的文档名列表"""
        return self.doc_names.copy()

    def get_stats(self) -> Dict[str, Any]:
        """返回知识库统计信息"""
        vector_store_size_mb = 0.0
        if os.path.exists(config.VECTOR_STORE_PATH):
            total_size = 0
            for dirpath, dirnames, filenames in os.walk(config.VECTOR_STORE_PATH):
                for f in filenames:
                    fp = os.path.join(dirpath, f)
                    if os.path.isfile(fp):
                        total_size += os.path.getsize(fp)
            vector_store_size_mb = total_size / (1024 * 1024)

        return {
            "num_docs": len(self.doc_names),
            "num_chunks": len(self.chunks),
            "vector_store_size_mb": vector_store_size_mb,
        }